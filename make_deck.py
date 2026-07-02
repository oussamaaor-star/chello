# -*- coding: utf-8 -*-
"""Génère la présentation Chello × AlAmeenSoft (.pptx) — style premium crème/ink."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Charte ──
INK       = RGBColor(0x18, 0x14, 0x0F)
CREAM     = RGBColor(0xFA, 0xF8, 0xF4)
CREAMD    = RGBColor(0xF1, 0xEC, 0xE2)
INKSOFT   = RGBColor(0x4A, 0x44, 0x3C)
ACCENT    = RGBColor(0xB9, 0xAB, 0x95)   # champagne/taupe (accent, PAS or vif)
SILVER    = RGBColor(0x9C, 0x94, 0x8A)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
HEAD = "Georgia"          # serif élégant (dispo Windows/Mac)
BODY = "Segoe UI"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5

def new(bg=CREAM):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = bg
    return s

def rect(s, l, t, w, h, color, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def rrect(s, l, t, w, h, color, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    try: sp.adjustments[0] = 0.06
    except Exception: pass
    return sp

def tb(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    b = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf

def para(tf, text, font=BODY, size=16, color=INKSOFT, bold=False, italic=False,
         align=PP_ALIGN.LEFT, after=6, line=1.15, first=False, upper=False):
    p = tf.paragraphs[0] if (first and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(after); p.space_before = Pt(0)
    try: p.line_spacing = line
    except Exception: pass
    r = p.add_run(); r.text = text.upper() if upper else text
    f = r.font; f.name = font; f.size = Pt(size); f.bold = bold; f.italic = italic; f.color.rgb = color
    return p

def eyebrow(s, text, color=SILVER, top=0.62, left=0.9):
    para(tb(s, left, top, 11, 0.4), text, BODY, 12, color, bold=True, upper=True, first=True)

def rule(s, l, t, w, color=ACCENT, h=0.028):
    rect(s, l, t, w, h, color)

def foot(s, n, dark=False):
    c = SILVER if not dark else RGBColor(0x8f,0x86,0x77)
    para(tb(s, 0.9, SH-0.5, 4, 0.3), "Chello", HEAD, 12, c, italic=True, first=True)
    para(tb(s, SW-1.4, SH-0.5, 0.9, 0.3), f"{n:02d}", BODY, 11, c, align=PP_ALIGN.RIGHT, first=True)

def title(s, text, top=1.15, left=0.9, size=40, color=INK, w=11.5):
    tf = tb(s, left, top, w, 1.9)
    for i, ln in enumerate(text.split("\n")):
        para(tf, ln, HEAD, size, color, italic=True, after=2, line=1.02, first=(i==0))
    return tf

def subtitle(s, text, top, left=0.9, color=INKSOFT, w=11.2, size=16):
    para(tb(s, left, top, w, 0.7), text, BODY, size, color, italic=True, first=True)

def bullets(s, items, top, left=0.9, w=11.4, size=17, gap=13, color=INKSOFT):
    tf = tb(s, left, top, w, SH-top-0.7)
    for i, (lead, rest) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.12
        d = p.add_run(); d.text = "—  "; d.font.name = BODY; d.font.size = Pt(size); d.font.color.rgb = ACCENT; d.font.bold = True
        if lead:
            a = p.add_run(); a.text = lead + " "; a.font.name = BODY; a.font.size = Pt(size); a.font.bold = True; a.font.color.rgb = INK
        b = p.add_run(); b.text = rest; b.font.name = BODY; b.font.size = Pt(size); b.font.color.rgb = color

# ══════════════════════════════════════════════════════════════ SLIDES

# 1 — COVER (ink)
s = new(INK)
rect(s, 0, 0, SW, 0.16, ACCENT)
para(tb(s, 1.1, 1.5, 11, 0.4), "Proposition · Chello × AlAmeenSoft", BODY, 13, ACCENT, bold=True, upper=True, first=True)
tf = tb(s, 1.1, 2.5, 11.2, 2.6)
para(tf, "4 000 produits.", HEAD, 60, CREAM, italic=True, after=2, line=1.0, first=True)
para(tf, "Zéro saisie à la main.", HEAD, 60, ACCENT, italic=True, after=2, line=1.0)
rule(s, 1.12, 5.25, 1.1, ACCENT)
para(tb(s, 1.1, 5.5, 10.5, 1.0), "La synchronisation automatique entre votre site et votre logiciel de gestion.",
     BODY, 18, RGBColor(0xd8,0xd2,0xc7), first=True)
para(tb(s, 1.1, SH-0.7, 11, 0.3), "Proposition préparée pour le propriétaire de Chello · Mascate, Oman", BODY, 11, SILVER, first=True)

# 2 — CONSTAT
s = new(CREAM); eyebrow(s, "Le constat")
title(s, "Votre boutique, c'est 4 000 produits."); rule(s, 0.92, 2.05, 1.0)
subtitle(s, "Une collection riche. Une vraie force. Mais aussi un vrai défi.", 2.25)
bullets(s, [
    ("", "4 000 références : noms, prix, photos, tailles, couleurs, stock."),
    ("", "Chaque jour le stock bouge : ventes en boutique, réassorts, retours."),
    ("", "Votre site doit refléter exactement ce que vous avez en rayon."),
    ("La vraie question :", "comment remplir et tenir à jour un catalogue de 4 000 produits, sans y passer vos journées ?"),
], 3.15)
foot(s, 2)

# 3 — PROBLÈME
s = new(CREAM); eyebrow(s, "Le problème caché")
title(s, "Saisir 4 000 produits à la main ? Intenable."); rule(s, 0.92, 2.05, 1.0)
subtitle(s, "Ce n'est pas un travail d'un jour. C'est un travail sans fin.", 2.25)
bullets(s, [
    ("Le temps :", "des semaines de saisie… puis à recommencer chaque jour pour le stock."),
    ("Les erreurs :", "un prix faux, une taille oubliée — sur 4 000 lignes, c'est inévitable."),
    ("La survente :", "vendu en ligne, mais déjà parti en boutique. Cliente déçue, remboursement."),
    ("Le coût caché :", "quelqu'un doit le faire… au lieu de vendre et de conseiller."),
], 3.15)
foot(s, 3)

# 4 — STATS (ink)
s = new(INK); eyebrow(s, "Le vrai coût du « à la main »", color=ACCENT)
title(s, "Ce que coûte la saisie manuelle.", color=CREAM); rule(s, 0.92, 2.05, 1.0, ACCENT)
stats = [("330+ h", "pour saisir 4 000 produits à la main"),
         ("≈ 2 mois", "de travail à plein temps, une personne"),
         ("1 000–2 500 €", "de main d'œuvre, juste pour démarrer")]
x = 0.9; cw = 3.9
for i, (big, cap) in enumerate(stats):
    lx = 0.9 + i*4.15
    para(tb(s, lx, 2.7, cw, 1.1), big, HEAD, 44, ACCENT, italic=True, first=True)
    rule(s, lx+0.02, 3.72, 0.7, RGBColor(0x5a,0x53,0x47))
    para(tb(s, lx, 3.9, cw, 1.1), cap, BODY, 15, RGBColor(0xcf,0xc8,0xbd), line=1.25, first=True)
para(tb(s, 0.9, 5.7, 11.4, 0.8),
     "…et il faudrait 2 à 5 h chaque jour pour tenir le stock à jour. Personne ne tient ce rythme — alors le stock en ligne devient faux en quelques jours.",
     BODY, 15, RGBColor(0xb9,0xab,0x95), italic=True, line=1.25, first=True)
foot(s, 4, dark=True)

# 5 — SOLUTION
s = new(CREAM); eyebrow(s, "La solution")
title(s, "On connecte votre site à AlAmeenSoft."); rule(s, 0.92, 2.05, 1.0)
subtitle(s, "Vos produits et votre stock remontent tout seuls. Vos commandes repartent toutes seules.", 2.25, w=11.4)
bullets(s, [
    ("", "AlAmeenSoft gère déjà vos produits, vos prix et votre stock."),
    ("", "On crée le pont entre votre logiciel et votre site."),
    ("Une seule saisie,", "dans AlAmeenSoft. Le site suit automatiquement."),
    ("", "Vous travaillez comme d'habitude — la technologie fait le reste."),
], 3.25)
foot(s, 5)

# 6 — AVANT / APRÈS (deux panneaux)
s = new(CREAM); eyebrow(s, "La différence")
title(s, "Le même site. Deux quotidiens opposés."); rule(s, 0.92, 2.05, 1.0)
av = ["Mise en ligne : un par un, à la main", "Stock : chaque jour, manuellement",
      "Erreurs : élevées", "Survente : fréquente", "Commandes : ressaisies à la main", "Votre temps : absorbé par la saisie"]
ap = ["Mise en ligne : automatique, d'un coup", "Stock : à jour en temps réel",
      "Erreurs : quasi nulles", "Survente : éliminée", "Commandes : envoyées automatiquement", "Votre temps : rendu à la vente"]
# panneau AVANT
rrect(s, 0.9, 2.55, 5.6, 4.3, CREAMD)
para(tb(s, 1.2, 2.8, 5, 0.5), "AVANT · saisie manuelle", BODY, 13, INKSOFT, bold=True, upper=True, first=True)
tfa = tb(s, 1.2, 3.4, 5.0, 3.3)
for i, t in enumerate(av): para(tfa, t, BODY, 14, INKSOFT, after=9, first=(i==0))
# panneau APRÈS
rrect(s, 6.85, 2.55, 5.6, 4.3, INK)
para(tb(s, 7.15, 2.8, 5, 0.5), "APRÈS · connecté", BODY, 13, ACCENT, bold=True, upper=True, first=True)
tfb = tb(s, 7.15, 3.4, 5.0, 3.3)
for i, t in enumerate(ap): para(tfb, t, BODY, 14, RGBColor(0xe9,0xe3,0xd9), after=9, first=(i==0))
foot(s, 6)

# 7 — BÉNÉFICES
s = new(CREAM); eyebrow(s, "Ce que vous y gagnez")
title(s, "Des résultats, pas de la technique."); rule(s, 0.92, 2.05, 1.0)
bullets(s, [
    ("Un stock toujours juste —", "en ligne = en rayon, sans y penser."),
    ("Zéro survente —", "vous ne vendez jamais ce que vous n'avez plus."),
    ("Zéro double saisie —", "vous saisissez une fois, dans AlAmeenSoft. Point."),
    ("Commandes → comptabilité —", "chaque vente en ligne rentre dans votre gestion."),
    ("Du temps rendu —", "vos équipes conseillent et vendent, au lieu de recopier."),
], 2.55, gap=15)
foot(s, 7)

# 8 — COMMENT ÇA MARCHE
s = new(CREAM); eyebrow(s, "Comment ça marche")
title(s, "Simple, en 3 étapes."); rule(s, 0.92, 2.05, 1.0)
steps = [("1", "On relie", "AlAmeenSoft à votre site, en coulisses."),
         ("2", "On synchronise", "les 4 000 produits + le stock, automatiquement."),
         ("3", "Ça tourne", "chaque vente repart dans l'ERP. En continu, sans vous.")]
for i, (n, h, d) in enumerate(steps):
    lx = 0.9 + i*4.15
    rrect(s, lx, 2.7, 3.8, 2.7, CREAMD)
    para(tb(s, lx+0.35, 2.95, 1, 0.9), n, HEAD, 40, ACCENT, italic=True, first=True)
    para(tb(s, lx+0.35, 3.85, 3.1, 0.5), h, HEAD, 20, INK, italic=True, first=True)
    para(tb(s, lx+0.35, 4.35, 3.1, 1.0), d, BODY, 14, INKSOFT, line=1.2, first=True)
para(tb(s, 0.9, 5.9, 11, 0.5), "Aucun changement dans votre façon de travailler. Votre logiciel reste votre logiciel — on lui donne une vitrine en ligne.",
     BODY, 14, SILVER, italic=True, first=True)
foot(s, 8)

# 9 — POURQUOI Y PENSER MAINTENANT
s = new(CREAM); eyebrow(s, "Pourquoi y penser maintenant")
title(s, "Mieux vaut l'intégrer dès la construction."); rule(s, 0.92, 2.05, 1.0)
subtitle(s, "L'anticiper pendant la création du site, c'est plus simple et plus solide que l'ajouter après coup.", 2.25, w=11.6)
bullets(s, [
    ("", "Le site est en cours de création — c'est le moment idéal pour prévoir la connexion."),
    ("", "Intégrée dès le départ, la synchronisation est plus propre, plus fiable, et sans mauvaise surprise."),
    ("", "Vous évitez d'avance de saisir 4 000 produits à la main — ou de lancer un site à moitié vide."),
    ("", "Le jour du lancement, votre catalogue est déjà complet et à jour, automatiquement."),
], 3.15)
foot(s, 9)

# 10 — LA DÉMARCHE (2 phases, sans prix)
s = new(CREAM); eyebrow(s, "La démarche")
title(s, "Une prestation claire, en 2 temps."); rule(s, 0.92, 2.05, 1.0)
phases = [("Phase 1 · Découverte", "On analyse AlAmeenSoft, on valide que la connexion est possible, et on définit le plan précis. Vous savez exactement où vous allez."),
          ("Phase 2 · Réalisation", "On construit le pont, on synchronise vos 4 000 produits et votre stock, et on teste avec vous jusqu'à ce que tout soit parfait.")]
for i, (h, d) in enumerate(phases):
    lx = 1.35 + i*5.35
    rrect(s, lx, 2.6, 4.9, 3.05, CREAMD)
    para(tb(s, lx+0.4, 2.95, 4.1, 0.7), h, HEAD, 21, INK, italic=True, line=1.05, first=True)
    rule(s, lx+0.42, 3.85, 0.6)
    para(tb(s, lx+0.4, 4.05, 4.1, 1.5), d, BODY, 14.5, INKSOFT, line=1.3, first=True)
para(tb(s, 1.2, 6.0, 10.9, 0.9),
     "Un module à part entière, qui apporte une vraie valeur à votre boutique — parlons-en ensemble pour définir ce qui vous convient.",
     BODY, 16, INK, italic=True, align=PP_ALIGN.CENTER, line=1.3, first=True)
foot(s, 10)

# 11 — OBJECTIONS
s = new(CREAM); eyebrow(s, "Vos questions, nos réponses")
title(s, "Trois doutes. Trois réponses."); rule(s, 0.92, 2.05, 1.0)
qa = [("« Est-ce compliqué pour moi ? »", "Non. Vous continuez à travailler dans AlAmeenSoft comme d'habitude. La connexion se fait en coulisses — vous n'avez rien à gérer."),
      ("« C'est risqué ? »", "On avance par phases. La Découverte valide que tout est possible avant d'engager la suite. Votre logiciel actuel n'est jamais mis en danger."),
      ("« Et si ça bug ? »", "Une fois la connexion en place, la synchronisation est automatique et fiable : pas d'erreur fonctionnelle à gérer de votre côté. Tout tourne seul.")]
tf = tb(s, 0.9, 2.55, 11.5, 4.3)
for i, (q, a) in enumerate(qa):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(6); r = p.add_run(); r.text = q
    r.font.name = HEAD; r.font.size = Pt(20); r.font.italic = True; r.font.color.rgb = INK
    p2 = tf.add_paragraph(); p2.space_after = Pt(18); p2.line_spacing = 1.2
    r2 = p2.add_run(); r2.text = a; r2.font.name = BODY; r2.font.size = Pt(15); r2.font.color.rgb = INKSOFT
foot(s, 11)

# 12 — CTA (ink)
s = new(INK)
rect(s, 0, SH-0.16, SW, 0.16, ACCENT)
para(tb(s, 1.1, 1.7, 11, 0.4), "Prochaine étape", BODY, 13, ACCENT, bold=True, upper=True, first=True)
tf = tb(s, 1.1, 2.6, 11.2, 2.2)
para(tf, "Prêt à voir votre boutique", HEAD, 46, CREAM, italic=True, after=2, line=1.03, first=True)
para(tf, "se remplir toute seule ?", HEAD, 46, ACCENT, italic=True, after=2, line=1.03)
para(tb(s, 1.1, 5.0, 10.5, 0.9), "On démarre par la Découverte — un premier pas simple, sans engagement lourd. On regarde AlAmeenSoft ensemble, et vous décidez en connaissance de cause.",
     BODY, 17, RGBColor(0xd8,0xd2,0xc7), line=1.3, first=True)
para(tb(s, 1.1, SH-0.9, 11, 0.4), "Oussama · [votre téléphone] · [votre email]", BODY, 14, SILVER, first=True)

out = r"C:\Users\RAFI\Desktop\Chello-Proposition-AlAmeenSoft.pptx"
prs.save(out)
print("OK ->", out, "|", len(prs.slides._sldIdLst), "slides")
